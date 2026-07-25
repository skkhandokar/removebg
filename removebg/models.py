import os
import subprocess
import tempfile
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from django.http import FileResponse
from pdf2docx import Converter
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path

class PdfConvertView(APIView):
    def post(self, request, *args, **kwargs):
        uploaded_file = request.FILES.get('file')
        conversion_type = request.data.get('conversion_type')

        if not uploaded_file or not conversion_type:
            return Response({'error': 'File and conversion_type are required.'}, status=status.HTTP_400_BAD_REQUEST)

        temp_dir = tempfile.mkdtemp()
        input_path = os.path.join(temp_dir, uploaded_file.name)

        with open(input_path, 'wb+') as destination:
            for chunk in uploaded_file.chunks():
                destination.write(chunk)

        try:
            # 1. PDF to Word
            if conversion_type == 'pdf-to-word':
                output_path = os.path.join(temp_dir, 'output.docx')
                cv = Converter(input_path)
                cv.convert(output_path, start=0, end=None)
                cv.close()
                return FileResponse(open(output_path, 'rb'), as_attachment=True, filename='converted.docx')

            # 2. PDF to Excel
            elif conversion_type == 'pdf-to-excel':
                output_path = os.path.join(temp_dir, 'output.xlsx')
                all_tables = []
                with pdfplumber.open(input_path) as pdf:
                    for page in pdf.pages:
                        extracted = page.extract_tables()
                        for table in extracted:
                            df = pd.DataFrame(table[1:], columns=table[0])
                            all_tables.append(df)
                
                if all_tables:
                    final_df = pd.concat(all_tables, ignore_index=True)
                    final_df.to_excel(output_path, index=False)
                else:
                    pd.DataFrame([{"Message": "No tables detected"}]).to_excel(output_path, index=False)
                
                return FileResponse(open(output_path, 'rb'), as_attachment=True, filename='converted.xlsx')

            # 3. PDF to PPT
            elif conversion_type == 'pdf-to-ppt':
                output_path = os.path.join(temp_dir, 'output.pptx')
                images = convert_from_path(input_path)
                prs = Presentation()
                blank_slide_layout = prs.slide_layouts[6]

                for i, img in enumerate(images):
                    img_path = os.path.join(temp_dir, f'page_{i}.png')
                    img.save(img_path, 'PNG')
                    slide = prs.slides.add_slide(blank_slide_layout)
                    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

                prs.save(output_path)
                return FileResponse(open(output_path, 'rb'), as_attachment=True, filename='converted.pptx')

            # 4. Office (Word/Excel/PPT) to PDF
            elif conversion_type == 'office-to-pdf':
                output_dir = temp_dir
                subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', input_path, '--outdir', output_dir], check=True)
                
                filename_without_ext = os.path.splitext(uploaded_file.name)[0]
                output_path = os.path.join(output_dir, f"{filename_without_ext}.pdf")
                return FileResponse(open(output_path, 'rb'), as_attachment=True, filename='converted.pdf')

            else:
                return Response({'error': 'Invalid conversion type.'}, status=status.HTTP_400_BAD_REQUEST)

        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)